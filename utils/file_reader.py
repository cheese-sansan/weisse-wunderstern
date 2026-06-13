"""
多格式文件读取模块

支持主流文档格式的文本内容提取。
原生支持（零依赖）：TXT, MD, JSON, CSV, HTML, XML, LOG, RST, TEX
可选依赖：PDF, DOCX, XLSX, PPTX, ODT/ODS/ODP, EPUB, RTF, 图片OCR
输出包含 raw_text、markdown_text、元信息、warnings 和 LaTeX 检测结果。
"""

import os
import json
import re


def read_file(file_path: str) -> dict:
    """
    读取文件并返回文本内容、Markdown 富文本和元信息。

    返回:
        dict: {
            "success": bool,
            "content": str,
            "markdown_text": str,
            "file_name": str,
            "file_type": str,
            "file_size": int,
            "error": str | None,
            "warnings": list[str],
            "contains_latex": bool,
        }
    """
    if not os.path.exists(file_path):
        return _error(file_path, f"文件不存在: {file_path}")

    if not os.path.isfile(file_path):
        return _error(file_path, f"路径不是文件: {file_path}")

    file_name = os.path.basename(file_path)
    _, ext = os.path.splitext(file_name)
    ext = ext.lower()
    file_size = os.path.getsize(file_path)

    # ── 纯文本类（Python 标准库，零依赖）──
    text_formats = {
        ".txt", ".md", ".markdown", ".csv", ".html", ".htm", ".xml",
        ".log", ".rst", ".tex", ".ini", ".cfg", ".conf", ".yaml", ".yml",
        ".toml", ".properties",
    }

    json_formats = {".json"}

    office_formats = {
        ".pdf":  _read_pdf,
        ".docx": _read_docx,
        ".xlsx": _read_xlsx,
        ".pptx": _read_pptx,
        ".odt":  _read_odt,
        ".ods":  _read_odt,
        ".odp":  _read_odt,
    }

    ebook_formats = {".epub": _read_epub}
    rtf_formats = {".rtf": _read_rtf}
    image_formats = {
        ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif", ".webp",
    }

    if ext in json_formats:
        handler = _read_json
    elif ext in office_formats:
        handler = office_formats[ext]
    elif ext in ebook_formats:
        handler = ebook_formats[ext]
    elif ext in rtf_formats:
        handler = rtf_formats[ext]
    elif ext in image_formats:
        handler = _read_image_ocr
    elif ext in text_formats:
        handler = _read_text
    else:
        handler = _read_text

    warnings = []
    try:
        result = handler(file_path)
        if isinstance(result, tuple):
            content, warnings = result
        else:
            content = result
        markdown_text = content
        contains_latex = _contains_latex(content)
        return {
            "success": True,
            "content": content,
            "markdown_text": markdown_text,
            "file_name": file_name,
            "file_type": ext,
            "file_size": file_size,
            "error": None,
            "warnings": warnings,
            "contains_latex": contains_latex,
        }
    except ImportError as e:
        return _error(file_path, str(e), warnings)
    except Exception as e:
        return _error(file_path, str(e), warnings)


def _error(file_path: str, message: str, warnings: list = None) -> dict:
    if os.path.exists(file_path) and ("/" in file_path or "\\" in file_path):
        file_name = os.path.basename(file_path)
    else:
        file_name = file_path
    return {
        "success": False,
        "content": "",
        "markdown_text": "",
        "file_name": file_name,
        "file_type": "",
        "file_size": 0,
        "error": message,
        "warnings": warnings or [],
        "contains_latex": False,
    }


# ═══════════════════════════════════════════════════════════════════
# 表格与公式检测
# ═══════════════════════════════════════════════════════════════════

def _contains_latex(text: str) -> bool:
    """检测文本是否包含 LaTeX 公式标记。"""
    patterns = [
        r"\$\$.+?\$\$",          # display math $$...$$
        r"\$(?!\$).+?\$(?!\$)",  # inline math $...$  (not $$)
        r"\\\[.+?\\\]",          # display math \[...\]
        r"\\\(.+?\\\)",          # inline math \(...\)
    ]
    for pat in patterns:
        if re.search(pat, text, re.DOTALL):
            return True
    return False


def _table_to_markdown(rows: list[list[str]]) -> str:
    """将二维表格数据转为 Markdown 表格字符串。"""
    if not rows or len(rows) < 2:
        return ""
    sanitized = [[str(c).replace("|", "&#124;").replace("\n", " ") for c in row] for row in rows]
    max_cols = max(len(r) for r in sanitized)
    for r in sanitized:
        while len(r) < max_cols:
            r.append("")
    lines = []
    lines.append("| " + " | ".join(sanitized[0]) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in sanitized[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════════
# 原生支持（Python 标准库，零依赖）
# ═══════════════════════════════════════════════════════════════════

def _read_text(file_path: str) -> str:
    """读取纯文本文件，自动检测编码（UTF-8 → GBK → latin-1）"""
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _read_json(file_path: str) -> str:
    """读取 JSON 文件并格式化输出"""
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)


# ═══════════════════════════════════════════════════════════════════
# PDF（可选依赖：pymupdf > pdfplumber > PyPDF2）
# ═══════════════════════════════════════════════════════════════════

def _read_pdf(file_path: str):
    """读取 PDF 并尝试提取表格为 Markdown。返回 (text, warnings)。"""
    warnings = []
    text = ""
    used_parser = None
    tables_md = []

    # 尝试 pymupdf
    try:
        import fitz
        doc = fitz.open(file_path)
        pages = [page.get_text() for page in doc]
        doc.close()
        text = "\n\n".join(pages)
        if text.strip():
            used_parser = "pymupdf"
    except ImportError:
        pass
    except Exception as e:
        warnings.append(f"pymupdf: {e}")

    # 尝试 pdfplumber（可提取表格）
    if not used_parser:
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        pages.append(t)
                    # 尝试提取表格
                    try:
                        tables = page.extract_tables()
                        for tbl in tables:
                            if tbl and len(tbl) > 1:
                                tables_md.append(_table_to_markdown(tbl))
                    except Exception:
                        pass
            text = "\n\n".join(pages)
            if text.strip():
                used_parser = "pdfplumber"
                if tables_md:
                    text += "\n\n---\n\n## 提取的表格\n\n" + "\n\n".join(tables_md)
        except ImportError:
            pass
        except Exception as e:
            warnings.append(f"pdfplumber: {e}")

    # 尝试 PyPDF2
    if not used_parser:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            pages = [page.extract_text() for page in reader.pages if page.extract_text()]
            text = "\n\n".join(pages)
            if text.strip():
                used_parser = "PyPDF2"
        except ImportError:
            pass
        except Exception as e:
            warnings.append(f"PyPDF2: {e}")

    if not used_parser:
        raise ImportError(
            "读取 PDF 需要安装第三方库。请执行以下任一命令：\n"
            "  pip install pymupdf\n"
            "  pip install pdfplumber\n"
            "  pip install PyPDF2"
        )

    if used_parser != "pymupdf" and not tables_md:
        warnings.append("PDF 表格提取未执行（仅 pdfplumber 支持 Markdown 表格）。安装 pdfplumber 可获得更佳表格提取。")

    return text, warnings


# ═══════════════════════════════════════════════════════════════════
# DOCX（可选依赖：python-docx）
# ═══════════════════════════════════════════════════════════════════

def _read_docx(file_path: str):
    """读取 DOCX 并将表格输出为 Markdown。返回 (text, warnings)。"""
    warnings = []
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []

        # 按文档顺序遍历 body 元素
        from docx.oxml.ns import qn
        body = doc.element.body
        para_idx = 0
        table_idx = 0

        for child in body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                if para_idx < len(doc.paragraphs):
                    p = doc.paragraphs[para_idx]
                    if p.text.strip():
                        parts.append(p.text.strip())
                    para_idx += 1
            elif tag == "tbl":
                if table_idx < len(doc.tables):
                    t = doc.tables[table_idx]
                    rows = []
                    for row in t.rows:
                        rows.append([cell.text for cell in row.cells])
                    if rows:
                        parts.append(_table_to_markdown(rows))
                    else:
                        for row in t.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    parts.append(cell.text.strip())
                    table_idx += 1

        # fallback: 如果顺序遍历未覆盖全部内容
        if not parts:
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())
            for t in doc.tables:
                rows = [[cell.text for cell in row.cells] for row in t.rows]
                if rows:
                    parts.append(_table_to_markdown(rows))

        return "\n\n".join(parts), warnings
    except ImportError:
        raise ImportError(
            "读取 DOCX 需要安装 python-docx。请执行：\n"
            "  pip install python-docx"
        )


# ═══════════════════════════════════════════════════════════════════
# XLSX（可选依赖：openpyxl）
# ═══════════════════════════════════════════════════════════════════

def _read_xlsx(file_path: str):
    """读取 XLSX 并转为 Markdown 表格。返回 (text, warnings)。"""
    warnings = []
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(list(row))
            if rows:
                parts.append(_table_to_markdown(rows))
        wb.close()
        return "\n\n".join(parts), warnings
    except ImportError:
        raise ImportError(
            "读取 XLSX 需要安装 openpyxl。请执行：\n"
            "  pip install openpyxl"
        )


# ═══════════════════════════════════════════════════════════════════
# PPTX（可选依赖：python-pptx）
# ═══════════════════════════════════════════════════════════════════

def _read_pptx(file_path: str):
    """读取 PPTX。返回 (text, warnings)。"""
    warnings = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_out = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            lines.append(para.text.strip())
                if shape.has_table:
                    table = shape.table
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    if rows:
                        lines.append(_table_to_markdown(rows))
            slides_out.append("\n".join(lines))
        return "\n\n".join(slides_out), warnings
    except ImportError:
        raise ImportError(
            "读取 PPTX 需要安装 python-pptx。请执行：\n"
            "  pip install python-pptx"
        )


# ═══════════════════════════════════════════════════════════════════
# ODT / ODS / ODP（可选依赖：odfpy）
# ═══════════════════════════════════════════════════════════════════

def _read_odt(file_path: str) -> str:
    try:
        from odf.opendocument import load
        from odf import text as odf_text
        doc = load(file_path)
        parts = []

        def _walk(element):
            if isinstance(element, odf_text.P):
                txt = ""
                for node in element.childNodes:
                    if node.nodeType == element.TEXT_NODE:
                        txt += node.data
                if txt.strip():
                    parts.append(txt.strip())
            for child in element.childNodes:
                _walk(child)

        _walk(doc.text)
        return "\n\n".join(parts)
    except ImportError:
        raise ImportError(
            "读取 ODT/ODS/ODP 需要安装 odfpy。请执行：\n"
            "  pip install odfpy"
        )


# ═══════════════════════════════════════════════════════════════════
# EPUB（可选依赖：ebooklib）
# ═══════════════════════════════════════════════════════════════════

def _read_epub(file_path: str) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from html.parser import HTMLParser

        class _Stripper(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts = []

            def handle_data(self, data):
                self.parts.append(data)

        book = epub.read_epub(file_path)
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                content = item.get_content().decode("utf-8", errors="replace")
                stripper = _Stripper()
                stripper.feed(content)
                text = "".join(stripper.parts).strip()
                if text:
                    chapters.append(text)

        if not chapters:
            return "[EPUB] 未提取到文本内容，可能为纯图片电子书"
        return "\n\n---\n\n".join(chapters)
    except ImportError:
        raise ImportError(
            "读取 EPUB 需要安装 ebooklib。请执行：\n"
            "  pip install ebooklib"
        )


# ═══════════════════════════════════════════════════════════════════
# RTF（可选依赖：striprtf）
# ═══════════════════════════════════════════════════════════════════

def _read_rtf(file_path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            rtf_content = f.read()
        return rtf_to_text(rtf_content)
    except ImportError:
        raise ImportError(
            "读取 RTF 需要安装 striprtf。请执行：\n"
            "  pip install striprtf"
        )


# ═══════════════════════════════════════════════════════════════════
# 图片 OCR（可选依赖：Pillow + pytesseract）
# ═══════════════════════════════════════════════════════════════════

def _read_image_ocr(file_path: str) -> str:
    try:
        from PIL import Image
        try:
            import pytesseract
        except ImportError:
            raise ImportError(
                "图片 OCR 需要安装 pytesseract。请执行：\n"
                "  pip install pytesseract Pillow\n"
                "  并安装 Tesseract 引擎：https://github.com/tesseract-ocr/tesseract"
            )
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        if not text.strip():
            return "[OCR] 图片中未检测到文字"
        return text.strip()
    except ImportError as e:
        raise ImportError(str(e))
